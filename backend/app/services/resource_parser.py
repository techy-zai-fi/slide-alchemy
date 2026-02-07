import re
import os
from pathlib import Path
from typing import Optional

from ..models.resource import ParsedContent, ResourceType


class ResourceParser:
    def parse_raw_text(self, text: str) -> ParsedContent:
        sections = self._extract_markdown_sections(text)
        return ParsedContent(
            text=text,
            sections=sections,
            metadata={
                "type": "raw_text",
                "char_count": len(text),
                "word_count": len(text.split()),
            },
        )

    def parse_markdown(self, content: str) -> ParsedContent:
        sections = self._extract_markdown_sections(content)
        clean_text = re.sub(r"[#*_`\[\]]", "", content).strip()
        return ParsedContent(
            text=clean_text,
            sections=sections,
            metadata={
                "type": "markdown",
                "char_count": len(content),
                "section_count": len(sections),
            },
        )

    def parse_txt(self, file_path: str) -> ParsedContent:
        text = Path(file_path).read_text(encoding="utf-8")
        return self.parse_raw_text(text)

    async def parse_pdf(self, file_path: str) -> ParsedContent:
        from PyPDF2 import PdfReader

        reader = PdfReader(file_path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append({"page": i + 1, "text": text})

        full_text = "\n\n".join(p["text"] for p in pages)
        return ParsedContent(
            text=full_text,
            sections=[{"title": f"Page {p['page']}", "content": p["text"]} for p in pages],
            metadata={
                "type": "pdf",
                "page_count": len(reader.pages),
                "char_count": len(full_text),
            },
        )

    async def parse_docx(self, file_path: str) -> ParsedContent:
        from docx import Document

        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n\n".join(paragraphs)

        sections = []
        current_section = None
        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                if current_section:
                    sections.append(current_section)
                current_section = {"title": para.text, "content": ""}
            elif current_section:
                current_section["content"] += para.text + "\n"
        if current_section:
            sections.append(current_section)

        return ParsedContent(
            text=full_text,
            sections=sections,
            metadata={"type": "docx", "paragraph_count": len(paragraphs), "char_count": len(full_text)},
        )

    async def parse_pptx(self, file_path: str) -> ParsedContent:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []
        sections = []

        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            slide_content = "\n".join(texts)
            slides_text.append(slide_content)
            sections.append({"title": f"Slide {i + 1}", "content": slide_content})

        full_text = "\n\n".join(slides_text)
        return ParsedContent(
            text=full_text,
            sections=sections,
            metadata={"type": "pptx", "slide_count": len(prs.slides), "char_count": len(full_text)},
        )

    async def parse_youtube(self, url: str) -> ParsedContent:
        from youtube_transcript_api import YouTubeTranscriptApi

        video_id = self._extract_youtube_id(url)
        if not video_id:
            return ParsedContent(text="", metadata={"error": "Invalid YouTube URL"})

        ytt_api = YouTubeTranscriptApi()
        transcript = ytt_api.fetch(video_id)
        segments = [
            {"start": seg.start, "text": seg.text}
            for seg in transcript
        ]
        full_text = " ".join(seg["text"] for seg in segments)

        return ParsedContent(
            text=full_text,
            sections=[{"title": f"Segment {i+1}", "content": s["text"]} for i, s in enumerate(segments[:20])],
            metadata={
                "type": "youtube",
                "video_id": video_id,
                "segment_count": len(segments),
                "char_count": len(full_text),
            },
        )

    async def parse_web_url(self, url: str) -> ParsedContent:
        import trafilatura

        downloaded = trafilatura.fetch_url(url)
        if not downloaded:
            return ParsedContent(text="", metadata={"error": f"Failed to fetch {url}"})

        text = trafilatura.extract(downloaded, include_links=False, include_images=False) or ""
        return ParsedContent(
            text=text,
            sections=self._extract_markdown_sections(text),
            metadata={"type": "web_url", "url": url, "char_count": len(text)},
        )

    async def parse_image(self, file_path: str) -> ParsedContent:
        try:
            import pytesseract
            from PIL import Image

            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return ParsedContent(
                text=text,
                metadata={"type": "image", "ocr": True, "char_count": len(text)},
            )
        except Exception as e:
            return ParsedContent(
                text="",
                metadata={"type": "image", "ocr": False, "error": str(e)},
            )

    async def parse(self, resource_type: ResourceType, source: str) -> ParsedContent:
        parsers = {
            ResourceType.RAW_TEXT: lambda s: self.parse_raw_text(s),
            ResourceType.MARKDOWN: lambda s: self.parse_markdown(s),
            ResourceType.TXT: lambda s: self.parse_txt(s),
        }
        async_parsers = {
            ResourceType.PDF: self.parse_pdf,
            ResourceType.DOCX: self.parse_docx,
            ResourceType.PPTX: self.parse_pptx,
            ResourceType.YOUTUBE: self.parse_youtube,
            ResourceType.WEB_URL: self.parse_web_url,
            ResourceType.IMAGE: self.parse_image,
        }

        if resource_type in parsers:
            return parsers[resource_type](source)
        if resource_type in async_parsers:
            return await async_parsers[resource_type](source)

        return ParsedContent(text=source, metadata={"type": "unknown"})

    def _extract_markdown_sections(self, text: str) -> list[dict]:
        sections = []
        current = None
        for line in text.split("\n"):
            heading_match = re.match(r"^(#{1,6})\s+(.+)", line)
            if heading_match:
                if current:
                    sections.append(current)
                current = {"title": heading_match.group(2), "content": ""}
            elif current:
                current["content"] += line + "\n"
        if current:
            sections.append(current)
        return sections

    def _extract_youtube_id(self, url: str) -> Optional[str]:
        patterns = [
            r"(?:youtube\.com/watch\?v=|youtu\.be/|youtube\.com/embed/)([a-zA-Z0-9_-]{11})",
        ]
        for pattern in patterns:
            match = re.search(pattern, url)
            if match:
                return match.group(1)
        return None
